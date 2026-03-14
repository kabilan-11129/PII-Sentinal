"""
file_parser.py — Extract plain text from uploaded files.

Supported formats:
  • TXT / LOG / MD  — read directly as plain text
  • CSV             — use pandas to join all cells
  • PDF             — use pdfminer.six
  • DOCX            — use python-docx
  • XLSX            — use pandas (openpyxl engine)
  • XLS             — use pandas (xlrd engine)
  • PPTX            — use python-pptx (all slide shapes)
  • EML             — use Python stdlib email module
  • JSON            — use stdlib json; extract all string values recursively
  • XML             — use stdlib xml.etree.ElementTree; extract all text nodes
  • HTML / HTM      — use stdlib html.parser; strip tags
  • RTF             — use striprtf if available; fallback to regex stripping
  • ZIP             — use Python stdlib zipfile (extract and scan supported members)
  • TAR             — use Python stdlib tarfile (.tar / .tar.gz / .tgz)
"""

import email as email_lib
import io
import json
import os
import re
import tarfile
import zipfile
import xml.etree.ElementTree as ET
from html.parser import HTMLParser

import docx
import pandas as pd
from pdfminer.high_level import extract_text as pdf_extract_text


# Maximum characters to extract per archive member (prevents memory exhaustion)
_ARCHIVE_MEMBER_LIMIT = 500_000

# File extensions that can be extracted from inside archives/emails
_INNER_EXTS = {
    ".txt", ".log", ".md",
    ".csv",
    ".pdf",
    ".docx", ".pptx", ".rtf",
    ".xlsx", ".xls",
    ".json", ".xml",
    ".html", ".htm",
    ".eml",
    ".odt", ".ods",
}


# ──────────────────────────────────────────────
# Individual format parsers
# ──────────────────────────────────────────────

def parse_txt(filepath: str) -> str:
    """Read a plain-text / log / markdown file and return its content."""
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def parse_csv(filepath: str) -> str:
    """Read a CSV file with pandas and concatenate all cell values."""
    try:
        df = pd.read_csv(filepath, dtype=str, keep_default_na=False)
        return "\n".join(" ".join(row.astype(str).values) for _, row in df.iterrows())
    except Exception as e:
        return f"[CSV parse error: {e}]"


def parse_pdf(filepath: str) -> str:
    """Extract text from a PDF using pdfminer.six."""
    try:
        return pdf_extract_text(filepath) or ""
    except Exception as e:
        return f"[PDF parse error: {e}]"


def parse_docx(filepath: str) -> str:
    """Extract paragraph text from a DOCX file using python-docx."""
    try:
        doc = docx.Document(filepath)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as e:
        return f"[DOCX parse error: {e}]"


def parse_xlsx(filepath: str) -> str:
    """Extract all cell values from an Excel workbook using pandas (openpyxl)."""
    try:
        xl = pd.ExcelFile(filepath, engine="openpyxl")
        parts = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet, dtype=str)
            df.fillna("", inplace=True)
            parts.append(f"[Sheet: {sheet}]")
            parts.extend(" ".join(row.astype(str).values) for _, row in df.iterrows())
        return "\n".join(parts)
    except Exception as e:
        return f"[XLSX parse error: {e}]"


def parse_xls(filepath: str) -> str:
    """Extract all cell values from a legacy .xls workbook using pandas (xlrd)."""
    try:
        xl = pd.ExcelFile(filepath, engine="xlrd")
        parts = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet, dtype=str)
            df.fillna("", inplace=True)
            parts.append(f"[Sheet: {sheet}]")
            parts.extend(" ".join(row.astype(str).values) for _, row in df.iterrows())
        return "\n".join(parts)
    except Exception as e:
        return f"[XLS parse error: {e}]"


def parse_pptx(filepath: str) -> str:
    """Extract text from all slides and shapes in a PowerPoint .pptx file."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            texts.append(f"[Slide {slide_num}]")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs if run.text.strip())
                        if line.strip():
                            texts.append(line)
                # Also extract table cells
                if shape.has_table:
                    for row in shape.table.rows:
                        texts.append(" | ".join(cell.text for cell in row.cells if cell.text.strip()))
        return "\n".join(texts)
    except ImportError:
        return "[PPTX parse error: python-pptx not installed. Run: pip install python-pptx]"
    except Exception as e:
        return f"[PPTX parse error: {e}]"


def parse_json(filepath: str) -> str:
    """Load a JSON file and recursively extract all string values."""
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            data = json.load(f)
        texts = []

        def _extract(obj):
            if isinstance(obj, str):
                texts.append(obj)
            elif isinstance(obj, dict):
                for v in obj.values():
                    _extract(v)
            elif isinstance(obj, list):
                for item in obj:
                    _extract(item)

        _extract(data)
        return "\n".join(texts)
    except Exception as e:
        return f"[JSON parse error: {e}]"


def parse_xml(filepath: str) -> str:
    """Parse an XML file and extract all text node content."""
    try:
        tree = ET.parse(filepath)
        root = tree.getroot()
        texts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                texts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                texts.append(elem.tail.strip())
            # Also include attribute values (can carry emails, names, etc.)
            for val in elem.attrib.values():
                if val.strip():
                    texts.append(val.strip())
        return "\n".join(texts)
    except Exception as e:
        return f"[XML parse error: {e}]"


class _HtmlTextExtractor(HTMLParser):
    """Minimal HTML parser that strips tags and skips script/style blocks."""

    def __init__(self):
        super().__init__()
        self._texts: list = []
        self._skip = False

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip = True
        # Extract href / content attribute values
        for attr, val in attrs:
            if attr in ("href", "content", "alt", "title", "placeholder") and val and val.strip():
                if not self._skip:
                    self._texts.append(val.strip())

    def handle_endtag(self, tag):
        if tag in ("script", "style"):
            self._skip = False

    def handle_data(self, data):
        if not self._skip and data.strip():
            self._texts.append(data.strip())

    def get_text(self) -> str:
        return "\n".join(self._texts)


def parse_html(filepath: str) -> str:
    """Extract visible text content from an HTML file."""
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        parser = _HtmlTextExtractor()
        parser.feed(content)
        return parser.get_text()
    except Exception as e:
        return f"[HTML parse error: {e}]"


def parse_rtf(filepath: str) -> str:
    """
    Extract plain text from an RTF file.
    Uses striprtf if available; falls back to basic regex stripping.
    """
    try:
        from striprtf.striprtf import rtf_to_text
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            return rtf_to_text(f.read())
    except ImportError:
        pass
    except Exception as e:
        return f"[RTF parse error: {e}]"

    # Fallback: basic RTF control-word stripping
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        content = re.sub(r"\\[a-z]+\-?\d* ?", " ", content)   # remove control words
        content = re.sub(r"\\\*[^;]+;", " ", content)          # remove destinations
        content = re.sub(r"[{}\\]", " ", content)              # remove braces/backslashes
        content = re.sub(r"\s{2,}", "\n", content)             # collapse whitespace
        return content.strip()
    except Exception as e:
        return f"[RTF parse error: {e}]"


def parse_eml(filepath: str) -> str:
    """
    Extract text from an .eml email file.
    Collects subject, from/to headers, and all text/* body parts.
    """
    try:
        with open(filepath, "rb") as f:
            msg = email_lib.message_from_bytes(f.read())

        parts = []
        for header in ("From", "To", "Cc", "Subject", "Reply-To"):
            val = msg.get(header, "")
            if val:
                parts.append(f"{header}: {val}")

        for part in msg.walk():
            ct = part.get_content_type()
            if ct in ("text/plain", "text/html"):
                payload = part.get_payload(decode=True)
                if payload:
                    text = payload.decode(part.get_content_charset() or "utf-8", errors="replace")
                    parts.append(text[:_ARCHIVE_MEMBER_LIMIT])

        return "\n".join(parts)
    except Exception as e:
        return f"[EML parse error: {e}]"


def _odf_bytes_to_text(data: bytes) -> str:
    """Extract text from ODT/ODS raw bytes (both are ZIP-compressed XML)."""
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            if "content.xml" not in zf.namelist():
                return ""
            xml_bytes = zf.read("content.xml")
        root = ET.fromstring(xml_bytes)
        texts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                texts.append(elem.text.strip())
            if elem.tail and elem.tail.strip():
                texts.append(elem.tail.strip())
        return "\n".join(texts)
    except Exception as e:
        return f"[ODF parse error: {e}]"


def parse_odt(filepath: str) -> str:
    """Extract text from an OpenDocument Text (.odt) file."""
    try:
        with open(filepath, "rb") as f:
            return _odf_bytes_to_text(f.read())
    except Exception as e:
        return f"[ODT parse error: {e}]"


def parse_ods(filepath: str) -> str:
    """Extract text from an OpenDocument Spreadsheet (.ods) file."""
    try:
        with open(filepath, "rb") as f:
            return _odf_bytes_to_text(f.read())
    except Exception as e:
        return f"[ODS parse error: {e}]"


def parse_msg(filepath: str) -> str:
    """
    Extract text from an Outlook .msg file.
    Uses extract-msg if installed; falls back to raw binary text extraction.
    """
    try:
        import extract_msg as _msg_lib
        msg = _msg_lib.openMsg(filepath)
        parts = []
        if getattr(msg, "subject", None):
            parts.append(f"Subject: {msg.subject}")
        if getattr(msg, "sender", None):
            parts.append(f"From: {msg.sender}")
        if getattr(msg, "to", None):
            parts.append(f"To: {msg.to}")
        if getattr(msg, "cc", None):
            parts.append(f"Cc: {msg.cc}")
        if getattr(msg, "body", None):
            parts.append(str(msg.body))
        msg.close()
        return "\n".join(parts)
    except ImportError:
        pass  # Fallback below
    except Exception as e:
        return f"[MSG parse error: {e}]"

    # Fallback: read the .msg as UTF-8 text (extracts some readable strings)
    try:
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            raw = f.read()
        # Strip non-printable chars
        import string
        printable = "".join(c for c in raw if c in string.printable)
        return printable
    except Exception as e:
        return f"[MSG parse error: {e}]"


def _text_from_bytes(data: bytes, ext: str) -> str:
    """Parse raw bytes as a known file type using an in-memory buffer."""
    try:
        buf = io.BytesIO(data)
        if ext in (".txt", ".log", ".md"):
            return data.decode("utf-8", errors="replace")
        elif ext == ".csv":
            df = pd.read_csv(buf, dtype=str, keep_default_na=False)
            return "\n".join(" ".join(r.astype(str).values) for _, r in df.iterrows())
        elif ext == ".pdf":
            return pdf_extract_text(buf) or ""
        elif ext == ".docx":
            doc = docx.Document(buf)
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext == ".xlsx":
            xl = pd.ExcelFile(buf, engine="openpyxl")
            rows = []
            for sheet in xl.sheet_names:
                df = xl.parse(sheet, dtype=str)
                df.fillna("", inplace=True)
                rows.extend(" ".join(r.astype(str).values) for _, r in df.iterrows())
            return "\n".join(rows)
        elif ext == ".xls":
            xl = pd.ExcelFile(buf, engine="xlrd")
            rows = []
            for sheet in xl.sheet_names:
                df = xl.parse(sheet, dtype=str)
                df.fillna("", inplace=True)
                rows.extend(" ".join(r.astype(str).values) for _, r in df.iterrows())
            return "\n".join(rows)
        elif ext == ".eml":
            msg = email_lib.message_from_bytes(data)
            parts = []
            for h in ("From", "To", "Cc", "Subject"):
                v = msg.get(h, "")
                if v:
                    parts.append(f"{h}: {v}")
            for part in msg.walk():
                if part.get_content_type() in ("text/plain", "text/html"):
                    pl = part.get_payload(decode=True)
                    if pl:
                        parts.append(pl.decode(part.get_content_charset() or "utf-8", errors="replace"))
            return "\n".join(parts)
        elif ext == ".json":
            obj = json.loads(data.decode("utf-8", errors="replace"))
            texts: list = []

            def _ext(o):
                if isinstance(o, str):
                    texts.append(o)
                elif isinstance(o, dict):
                    for v in o.values():
                        _ext(v)
                elif isinstance(o, list):
                    for i in o:
                        _ext(i)

            _ext(obj)
            return "\n".join(texts)
        elif ext == ".xml":
            root = ET.fromstring(data.decode("utf-8", errors="replace"))
            texts = []
            for elem in root.iter():
                if elem.text and elem.text.strip():
                    texts.append(elem.text.strip())
                if elem.tail and elem.tail.strip():
                    texts.append(elem.tail.strip())
            return "\n".join(texts)
        elif ext in (".html", ".htm"):
            content = data.decode("utf-8", errors="replace")
            parser = _HtmlTextExtractor()
            parser.feed(content)
            return parser.get_text()
        elif ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(buf)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                t = " ".join(r.text for r in para.runs if r.text.strip())
                                if t.strip():
                                    texts.append(t)
                return "\n".join(texts)
            except ImportError:
                return ""
        elif ext == ".rtf":
            try:
                from striprtf.striprtf import rtf_to_text
                return rtf_to_text(data.decode("utf-8", errors="replace"))
            except ImportError:
                # Basic fallback stripping
                content = data.decode("utf-8", errors="replace")
                content = re.sub(r"\\[a-z]+\-?\d* ?", " ", content)
                content = re.sub(r"[{}\\]", " ", content)
                return re.sub(r"\s{2,}", "\n", content).strip()
        elif ext in (".odt", ".ods"):
            return _odf_bytes_to_text(data)
    except Exception:
        pass
    return ""


def parse_zip(filepath: str) -> str:
    """
    Extract and concatenate text from all supported files inside a ZIP archive.
    """
    parts = []
    try:
        with zipfile.ZipFile(filepath, "r") as zf:
            for info in zf.infolist():
                if info.is_dir():
                    continue
                ext = os.path.splitext(info.filename)[1].lower()
                if ext not in _INNER_EXTS:
                    continue
                try:
                    data = zf.read(info.filename)[:_ARCHIVE_MEMBER_LIMIT]
                    text = _text_from_bytes(data, ext)
                    if text:
                        parts.append(f"[{info.filename}]\n{text}")
                except Exception:
                    continue
    except Exception as e:
        return f"[ZIP parse error: {e}]"
    return "\n\n".join(parts) if parts else "[ZIP: no supported files found]"


def parse_tar(filepath: str) -> str:
    """
    Extract and concatenate text from all supported files inside a TAR archive
    (supports .tar, .tar.gz, .tgz).
    """
    parts = []
    try:
        with tarfile.open(filepath, "r:*") as tf:
            for member in tf.getmembers():
                if not member.isfile():
                    continue
                ext = os.path.splitext(member.name)[1].lower()
                if ext not in _INNER_EXTS:
                    continue
                try:
                    f = tf.extractfile(member)
                    if f is None:
                        continue
                    data = f.read(_ARCHIVE_MEMBER_LIMIT)
                    text = _text_from_bytes(data, ext)
                    if text:
                        parts.append(f"[{member.name}]\n{text}")
                except Exception:
                    continue
    except Exception as e:
        return f"[TAR parse error: {e}]"
    return "\n\n".join(parts) if parts else "[TAR: no supported files found]"


# ──────────────────────────────────────────────
# Dispatcher
# ──────────────────────────────────────────────

_PARSERS = {
    # Plain text variants
    ".txt"  : parse_txt,
    ".log"  : parse_txt,
    ".md"   : parse_txt,
    # Tabular data
    ".csv"  : parse_csv,
    ".xlsx" : parse_xlsx,
    ".xls"  : parse_xls,
    ".ods"  : parse_ods,
    # Documents
    ".pdf"  : parse_pdf,
    ".docx" : parse_docx,
    ".pptx" : parse_pptx,
    ".rtf"  : parse_rtf,
    ".odt"  : parse_odt,
    # Structured data
    ".json" : parse_json,
    ".xml"  : parse_xml,
    # Web
    ".html" : parse_html,
    ".htm"  : parse_html,
    # Email
    ".eml"  : parse_eml,
    ".msg"  : parse_msg,
    # Archives
    ".zip"  : parse_zip,
    ".tar"  : parse_tar,
    ".gz"   : parse_tar,   # handles .tar.gz; plain .gz is uncommon
    ".tgz"  : parse_tar,
}


def parse_file(filepath: str) -> str:
    """
    Auto-detect the file type from its extension and extract text.
    Returns the extracted text or an error message string.
    """
    ext = os.path.splitext(filepath)[1].lower()
    parser = _PARSERS.get(ext)
    if parser is None:
        return f"[Unsupported file type: {ext}]"
    return parser(filepath)
